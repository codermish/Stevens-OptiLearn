import openai
import streamlit as st
import json
from openai import OpenAI
from datetime import datetime
from pymongo import MongoClient
import re
import os
import numpy as np
from sklearn.metrics.pairwise import cosine_similarity
from io import BytesIO
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas

from io import BytesIO
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from reportlab.lib.utils import simpleSplit
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.units import inch

# Page setup
st.set_page_config(page_title="Stevens OptiLearn", layout="wide", initial_sidebar_state="collapsed")

# Enhanced CSS Styling with Stevens Institute of Technology branding
st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Roboto:wght@300;400;500;700;900&family=Roboto+Slab:wght@300;400;500;700&display=swap');

            
/* Global Styles */


html, body, [class*="css"] { 
    font-family: 'Roboto', -apple-system, BlinkMacSystemFont, 'Segoe UI', sans-serif; 
    color: #212529;
    line-height: 1.6;
}

.main .block-container { 
    padding: 1rem !important; 
    max-width: none !important; 
    margin: 0 !important; 
}

/* Prevent flickering */
div[data-testid="stAppViewContainer"] { transition: none !important; }

/* Header Section - Stevens Branding */
.header-section { 
    background: linear-gradient(135deg,#800020 0%, #a61e42 100%); 
    padding: 3.5rem 2rem;
    text-align: center;
    color: white;
    margin: 0;
    border-radius: 0;
    box-shadow: 0 4px 20px rgba(0, 0, 0, 0.15);
    position: relative;
    overflow: hidden;
}

.header-section::before {
    content: '';
    position: absolute;
    top: -50%;
    right: -10%;
    width: 60%;
    height: 200%;
    background: rgba(255, 255, 255, 0.05);
    transform: rotate(45deg);
}

.header-title { font-family: 'Crimson Text', serif; font-size: 3rem; font-weight: 700; margin: 0; color: white; }
.header-subtitle { font-size: 1.1rem; font-weight: 300; margin-top: 0.5rem; opacity: 0.95; color: white;}

/* Form Wrapper */
.form-wrapper { 
    max-width: 1000px;
    margin: 2rem auto;
    padding: 3rem;
    background: white;
    border-radius: 12px;
    box-shadow: 0 10px 40px rgba(0, 0, 0, 0.08);
    border: 1px solid rgba(0, 0, 0, 0.05);
    position: relative;
}

.form-wrapper::before {
    content: '';
    position: absolute;
    top: 0;
    left: 0;
    right: 0;
    height: 4px;
    background: linear-gradient(90deg, #C8102E 0%, #F7941E 50%, #C8102E 100%);
    border-radius: 12px 12px 0 0;
}

.section-title { 
    font-family: 'Roboto Slab', serif;
    font-size: 2rem;
    font-weight: 600;
    color: #800020;
    margin-bottom: 2rem;
    padding-bottom: 1rem;
    border-bottom: 3px solid #f0f0f0;
    position: relative;
}

.section-title::after {
    content: '';
    position: absolute;
    bottom: -3px;
    left: 0;
    width: 60px;
    height: 3px;
    background: #C8102E;
}

/* AI Chat Container */
.ai-chat-container { 
    background: white;
    border-radius: 16px;
    padding: 2.5rem;
    margin: 2rem auto;
    max-width: 1000px;
    box-shadow: 0 10px 40px rgba(0, 0, 0, 0.08);
    border: 1px solid #e9ecef;
}

.ai-message { 
    background: #f8f9fa;
    padding: 1.75rem;
    border-radius: 12px;
    border-left: 4px solid #C8102E;
    margin-bottom: 1.5rem;
    box-shadow: 0 2px 12px rgba(0, 0, 0, 0.04);
    position: relative;
    transition: all 0s ease;
}

.ai-message:hover {
    transform: translateY(-2px);
    box-shadow: 0 4px 20px rgba(0, 0, 0, 0.08);
}



/* User Message */
.user-message { 
    background: linear-gradient(135deg, #C8102E 0%, #9B0E23 100%);
    color: white;
    padding: 1.25rem 1.75rem;
    border-radius: 12px;
    margin-bottom: 1.5rem;
    margin-left: 3rem;
    box-shadow: 0 4px 16px rgba(200, 16, 46, 0.2);
    position: relative;
}

.ai-label { 
    font-size: 0.95rem;
    font-weight: 600;
    color: #C8102E;
    margin-bottom: 0.75rem;
    display: flex;
    align-items: center;
    gap: 0.75rem;
    text-transform: uppercase;
    letter-spacing: 0.5px;
}

.ai-label::before { 
    content: "🎓";
    font-size: 1.4rem;
}

/* Ratio Display */
.ratio-display { 
    display: flex;
    justify-content: space-around;
    margin: 3rem 0;
    gap: 2rem;
}

.ratio-item { 
    flex: 1;
    text-align: center;
    padding: 2rem;
    background: white;
    border-radius: 12px;
    border: 2px solid #e9ecef;
    box-shadow: 0 4px 20px rgba(0, 0, 0, 0.06);
    transition: all 0s ease;
    position: relative;
    overflow: hidden;
}

.ratio-item:hover {
    transform: translateY(-4px);
    box-shadow: 0 8px 30px rgba(0, 0, 0, 0.1);
    border-color: #C8102E;
}

.ratio-value { 
    font-size: 3rem;
    font-weight: 700;
    color: #C8102E;
    font-family: 'Roboto Slab', serif;
}

.ratio-label { 
    font-size: 1.1rem;
    font-weight: 500;
    color: #495057;
    margin-top: 0.5rem;
    text-transform: uppercase;
    letter-spacing: 1px;
}

/* Progress Indicator */
.progress-indicator { 
    display: flex;
    align-items: center;
    justify-content: center;
    gap: 1rem;
    margin: 2rem auto;
    padding: 1.5rem;
    background: white;
    border-radius: 12px;
    box-shadow: 0 4px 20px rgba(0, 0, 0, 0.06);
    max-width: 1000px;
    flex-wrap: wrap;
}

.progress-step { 
    padding: 0.75rem 1.5rem;
    background: #e9ecef;
    border-radius: 25px;
    font-size: 0.95rem;
    font-weight: 500;
    color: #6c757d;
    transition: all 0s ease;
    border: 2px solid transparent;
}

.progress-step.active { 
    background: linear-gradient(135deg,#800020 0%, #a61e42 100%); 
    color: white;
    transform: scale(1.05);
    box-shadow: 0 4px 16px rgba(200, 16, 46, 0.3);
}

.progress-step.completed { 
    background: #28a745;
    color: white;
    border-color: #28a745;
}

/* Professional Recommendation Cards */
.recommendations-container { 
    max-width: 1400px;
    margin: 2rem auto;
    padding: 0 1rem;
}

.recommendations-header { 
    text-align: center;
    margin-bottom: 3rem;
    padding: 3rem 2rem;
    background: linear-gradient(135deg, #C8102E 0%, #9B0E23 100%);
    color: white;
    border-radius: 16px;
    box-shadow: 0 12px 40px rgba(200, 16, 46, 0.25);
    position: relative;
    overflow: hidden;
}

.recommendations-header::before {
    content: '';
    position: absolute;
    top: -50%;
    left: -10%;
    width: 120%;
    height: 200%;
    background: radial-gradient(circle, rgba(255,255,255,0.1) 0%, transparent 70%);
}

.recommendations-header h1 { 
    font-family: 'Roboto Slab', serif;
    font-size: 3rem;
    color: white;
    margin-bottom: 1rem;
    text-shadow: 2px 2px 4px rgba(0,0,0,0.2);
    position: relative;
    z-index: 1;
}

.recommendations-header p { 
    font-size: 1.3rem;
    color: rgba(255,255,255,0.95);
    font-weight: 300;
    position: relative;
    z-index: 1;
}

/* Professional Course Card */
.professional-course-card {
    background: white;
    border-radius: 16px;
    box-shadow: 0 6px 30px rgba(0,0,0,0.08);
    border: 1px solid #e9ecef;
    overflow: hidden;
    transition: all 0s cubic-bezier(0.165, 0.84, 0.44, 1);
    position: relative;
    margin-bottom: 2rem;
}

.professional-course-card:hover {
    transform: translateY(-6px);
    box-shadow: 0 12px 50px rgba(200, 16, 46, 0.15);
    border-color: #C8102E;
}

.professional-course-card::before {
    content: '';
    position: absolute;
    top: 0;
    left: 0;
    right: 0;
    height: 5px;
    background: linear-gradient(90deg, #C8102E 0%, #F7941E 50%, #C8102E 100%);
}

.course-header-pro {
    background: linear-gradient(135deg, #2c3e50 0%, #34495e 100%);
    color: white;
    padding: 2.5rem 2rem;
    position: relative;
}

.course-rank-pro {
    position: absolute;
    top: 1.5rem;
    right: 1.5rem;
    background: linear-gradient(135deg, #C8102E 0%, #9B0E23 100%);
    color: white;
    padding: 0.5rem 1.25rem;
    border-radius: 25px;
    font-size: 0.95rem;
    font-weight: 700;
    box-shadow: 0 4px 12px rgba(200, 16, 46, 0.4);
}

.course-code-pro {
    font-size: 1rem;
    opacity: 0.85;
    margin-bottom: 0.5rem;
    font-weight: 500;
    letter-spacing: 1px;
    text-transform: uppercase;
}

.course-name-pro {
    font-size: 1.6rem;
    font-weight: 700;
    line-height: 1.3;
    margin: 0;
    font-family: 'Roboto Slab', serif;
    color: white;
}

.course-type-badge {
    display: inline-flex;
    align-items: center;
    gap: 0.5rem;
    background: rgba(255,255,255,0.2);
    backdrop-filter: blur(10px);
    padding: 0.5rem 1rem;
    border-radius: 20px;
    font-size: 0.9rem;
    font-weight: 600;
    margin-top: 1rem;
    border: 1px solid rgba(255,255,255,0.3);
}

.course-body-pro {
    padding: 2.5rem;
}

.similarity-section-pro {
    display: grid;
    grid-template-columns: 1fr 1fr;
    gap: 1.5rem;
    margin-bottom: 2rem;
}

.similarity-metric {
    background: #f8f9fa;
    padding: 1.5rem;
    border-radius: 12px;
    text-align: center;
    border: 1px solid #e9ecef;
    position: relative;
    transition: all 0s ease;
}

.similarity-metric:hover {
    background: #fff;
    box-shadow: 0 4px 16px rgba(0,0,0,0.06);
}

.similarity-metric::before {
    content: '';
    position: absolute;
    top: 0;
    left: 0;
    right: 0;
    height: 3px;
    background: linear-gradient(90deg, #C8102E 0%, #F7941E 100%);
    border-radius: 12px 12px 0 0;
}

.similarity-value-pro {
    font-size: 2rem;
    font-weight: 800;
    color: #C8102E;
    margin-bottom: 0.25rem;
    font-family: 'Roboto Slab', serif;
}

.similarity-label-pro {
    font-size: 0.9rem;
    color: #6c757d;
    font-weight: 500;
    text-transform: uppercase;
    letter-spacing: 0.5px;
}

.course-description-pro {
    background: #f8f9fa;
    padding: 1.75rem;
    border-radius: 12px;
    margin-bottom: 1.5rem;
    font-size: 1rem;
    color: #495057;
    line-height: 1.7;
    border-left: 4px solid #C8102E;
}

/* Recommendations Summary */
.recommendations-summary {
    background: linear-gradient(135deg,#800020 0%, #a61e42 100%); 
    color: white;
    padding: 3rem;
    border-radius: 16px;
    margin: 3rem auto;
    box-shadow: 0 12px 50px rgba(0,0,0,0.15);
    max-width: 1200px;
}

.summary-grid {
    display: grid;
    grid-template-columns: repeat(auto-fit, minmax(180px, 1fr));
    gap: 2rem;
    margin-top: 2rem;
}

.summary-metric {
    text-align: center;
    padding: 1.5rem;
    background: #808080;
    border-radius: 12px;
    backdrop-filter: blur(10px);
    border: 1px solid rgba(255,255,255,0.15);
    transition: all 0 ease;
}

.summary-metric:hover {
    background: #808080;
    transform: translateY(-2px);
}

.summary-metric-value {
    font-size: 2.5rem;
    font-weight: 800;
    color: white;
    margin-bottom: 0.5rem;
    font-family: 'Roboto Slab', serif;
}

.summary-metric-label {
    font-size: 0.95rem;
    color: rgba(255,255,255,0.9);
    text-transform: uppercase;
    letter-spacing: 0.5px;
    font-weight: 500;
}

/* Action Section */
.action-section {
    background: white;
    padding: 3rem;
    border-radius: 16px;
    margin: 3rem auto;
    max-width: 1000px;
    border: 1px solid #e9ecef;
    box-shadow: 0 6px 30px rgba(0,0,0,0.08);
}

.action-section h3 {
    color: #C8102E;
    font-family: 'Roboto Slab', serif;
    font-size: 1.8rem;
    margin-bottom: 2rem;
}

/* Email Error */
.email-error {
    color: #dc3545;
    font-size: 0.9rem;
    margin-top: 0.5rem;
    display: flex;
    align-items: center;
    gap: 0.5rem;
}

/* Buttons */
.stButton > button { 
    background: linear-gradient(135deg,#800020 0%, #a61e42 100%);  !important;
    color: white !important;
    font-weight: 600 !important;
    padding: 0.875rem 2.5rem !important;
    border: none !important;
    border-radius: 30px !important;
    transition: all 0s ease !important;
    font-size: 1rem !important;
    letter-spacing: 0.5px !important;
    text-transform: uppercase !important;
    box-shadow: 0 4px 16px rgba(200, 16, 46, 0.2) !important;
}

.stButton > button:hover { 
    background: linear-gradient(135deg,#800020 0%, #a61e42 100%);  !important;
    transform: translateY(-2px) !important;
    box-shadow: 0 6px 24px rgba(200, 16, 46, 0.3) !important;
    color: white;
}
/* make sure disabled buttons stay white text (and bump opacity if you like) */
.stButton > button:disabled {
  color: white !important;
  opacity: 1 !important;
}


/* Input Fields */
.stTextInput > div > div > input,
.stSelectbox > div > div > select,
.stTextArea > div > div > textarea {
    border: 2px solid #e9ecef !important;
    border-radius: 8px !important;
    padding: 0.75rem !important;
    font-size: 1rem !important;
    transition: all 0s ease !important;
}

.stTextInput > div > div > input:focus,
.stSelectbox > div > div > select:focus,
.stTextArea > div > div > textarea:focus {
    border-color: #C8102E !important;
    box-shadow: 0 0 0 0.2rem rgba(200, 16, 46, 0.15) !important;
}

/* Slider */
.stSlider > div > div > div {
    background: #C8102E !important;
}

.stSlider > div > div > div > div {
    background: #C8102E !important;
    box-shadow: 0 2px 8px rgba(200, 16, 46, 0.3) !important;
}

/* Hide Streamlit elements */
#MainMenu {visibility: hidden;} 
footer {visibility: hidden;} 
header {visibility: hidden;}
.stDeployButton {display: none;}

/* Responsive Design */
@media (max-width: 768px) {
    .header-title { font-size: 2.5rem; }
    .header-subtitle { font-size: 1rem; }
    .form-wrapper { padding: 2rem 1.5rem; }
    .section-title { font-size: 1.6rem; }
    .progress-indicator { gap: 0.5rem; }
    .progress-step { padding: 0.5rem 1rem; font-size: 0.85rem; }
    .ratio-display { flex-direction: column; gap: 1rem; }
    .similarity-section-pro { grid-template-columns: 1fr; }
    .summary-grid { grid-template-columns: repeat(2, 1fr); gap: 1rem; }
}

/* Professional Typography */
h1, h2, h3, h4, h5, h6 {
    font-family: 'Roboto Slab', serif;
    color: #212529;
}

p {
    line-height: 1.8;
    color: #495057;
}

/* No Results State */
.no-results {
    text-align: center;
    padding: 3rem;
    background: #fff5f5;
    border-radius: 16px;
    border: 2px solid #ffdddd;
    margin: 2rem auto;
    max-width: 600px;
}

.no-results h3 {
    color: #C8102E;
    font-family: 'Roboto Slab', serif;
    margin-bottom: 1rem;
}

.no-results p {
    color: #666;
    line-height: 1.6;
}
            

.stApp {
  background: linear-gradient(135deg, #f8f9fa 0%, #e9ecef 100%) fixed no-repeat !important;
  background-attachment: fixed !important;
}

</style>
""", unsafe_allow_html=True)

# Initialize session state
if 'logged_in' not in st.session_state: st.session_state.logged_in = False
if 'current_step' not in st.session_state: st.session_state.current_step = 1
if 'form_data' not in st.session_state: st.session_state.form_data = {}
if 'initial_req' not in st.session_state: st.session_state.initial_req = None
if 'conversation_history' not in st.session_state: st.session_state.conversation_history = []
if 'question_count' not in st.session_state: st.session_state.question_count = 0
if 'requirements_data' not in st.session_state: st.session_state.requirements_data = {}
if 'ai_analysis' not in st.session_state: st.session_state.ai_analysis = {}
if 'openai_api_key' not in st.session_state: st.session_state.openai_api_key = ""
if 'course_recommendations' not in st.session_state: st.session_state.course_recommendations = []
if 'current_ai_question' not in st.session_state: st.session_state.current_ai_question = None

USER_CREDENTIALS = {"stevens_org": "admin", "admin": "admin"}

# Email validation function
def is_valid_email(email):
    pattern = r'^[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}$'
    return re.match(pattern, email) is not None

def load_course_embeddings():
    """Load course embeddings from JSON file"""
    try:
        # Try to read the embeddings file
        if 'all_course_embeddings.json' in [f for f in os.listdir('.') if f.endswith('.json')]:
            with open('all_course_embeddings.json', 'r') as f:
                embeddings_data = json.load(f)
                return embeddings_data
        else:
            st.warning("Course embeddings file not found. Using sample data.")
            return []
    except Exception as e:
        st.error(f"Error loading course embeddings: {str(e)}")
        return []
    

# def generate_pdf(recommendations):
#     buffer = BytesIO()
#     c = canvas.Canvas(buffer, pagesize=letter)
#     width, height = letter
#     y = height - 50
#     margin_left = 50
#     margin_right = width - 50
#     line_height = 15
    
#     def check_page_break(current_y, needed_space=60):
#         """Check if we need a new page and create one if necessary"""
#         if current_y < needed_space:
#             c.showPage()
#             return height - 50
#         return current_y
    
#     def draw_wrapped_text(text, x, y, max_width, font_name="Helvetica", font_size=11):
#         """Draw text with proper wrapping"""
#         c.setFont(font_name, font_size)
        
#         # Split text into lines that fit within max_width
#         lines = simpleSplit(text, font_name, font_size, max_width)
        
#         current_y = y
#         for line in lines:
#             current_y = check_page_break(current_y)
#             c.drawString(x, current_y, line)
#             current_y -= line_height
        
#         return current_y
    
#     # Title
#     c.setFont("Helvetica-Bold", 18)
#     c.drawString(margin_left, y, "Course Recommendations Report")
#     y -= 40
    
#     # Add generation date
#     from datetime import datetime
#     c.setFont("Helvetica", 10)
#     c.drawString(margin_left, y, f"Generated on: {datetime.now().strftime('%B %d, %Y at %I:%M %p')}")
#     y -= 30
    
#     # Draw a line separator
#     c.line(margin_left, y, margin_right, y)
#     y -= 20
    
#     for i, course in enumerate(recommendations, 1):
#         y = check_page_break(y, 100)  # Ensure we have space for at least the header
        
#         # Extract course information
#         code = course.get('course_code', 'N/A')
#         name = course.get('course_name', 'Untitled Course')
#         desc = course.get('description', 'No description available.')
#         outcomes = course.get('learning_outcomes', []) or course.get('outcomes', [])
        
#         # Handle similarity score
#         similarity = course.get('adjusted_similarity', 0)
#         if hasattr(similarity, '__call__'):  # If it's a function reference
#             sim_text = "N/A"
#         else:
#             try:
#                 sim_percentage = convert_similarity_to_percentage(similarity) if 'convert_similarity_to_percentage' in globals() else f"{float(similarity)*100:.1f}"
#                 sim_text = f"{sim_percentage}%"
#             except:
#                 sim_text = "N/A"
        
#         # Course header with background
#         c.setFillColorRGB(0.9, 0.9, 0.9)  # Light gray background
#         c.rect(margin_left-5, y-5, margin_right-margin_left+10, 25, fill=1, stroke=0)
#         c.setFillColorRGB(0, 0, 0)  # Reset to black text
        
#         c.setFont("Helvetica-Bold", 14)
#         header_text = f"{i}. {code} - {name}"
#         if len(header_text) > 70:  # Truncate if too long
#             header_text = header_text[:67] + "..."
#         c.drawString(margin_left, y+5, header_text)
        
#         # Match percentage
#         c.setFont("Helvetica-Bold", 12)
#         c.setFillColorRGB(0, 0.5, 0)  # Green color for match percentage
#         c.drawString(margin_right - 80, y+5, f"Match: {sim_text}")
#         c.setFillColorRGB(0, 0, 0)  # Reset to black
        
#         y -= 35
        
#         # Description section
#         c.setFont("Helvetica-Bold", 11)
#         y = check_page_break(y)
#         c.drawString(margin_left + 10, y, "Description:")
#         y -= 18
        
#         # Description content with wrapping
#         max_desc_width = margin_right - margin_left - 20
#         y = draw_wrapped_text(desc, margin_left + 20, y, max_desc_width, "Helvetica", 10)
#         y -= 10
        
#         # Learning Outcomes section
#         if outcomes and len(outcomes) > 0:
#             y = check_page_break(y)
#             c.setFont("Helvetica-Bold", 11)
#             c.drawString(margin_left + 10, y, "Learning Outcomes:")
#             y -= 18
            
#             for outcome in outcomes:
#                 y = check_page_break(y)
                
#                 # Clean up the outcome text
#                 outcome_text = str(outcome).strip()
#                 if outcome_text:
#                     # Draw bullet point
#                     c.setFont("Helvetica", 10)
#                     c.drawString(margin_left + 20, y, "•")
                    
#                     # Draw outcome with wrapping
#                     max_outcome_width = margin_right - margin_left - 40
#                     y = draw_wrapped_text(outcome_text, margin_left + 35, y, max_outcome_width, "Helvetica", 10)
#                     y -= 5
        
#         # Add spacing between courses
#         y -= 20
        
#         # Draw separator line between courses (except for the last one)
#         if i < len(recommendations):
#             y = check_page_break(y)
#             c.setStrokeColorRGB(0.8, 0.8, 0.8)
#             c.line(margin_left, y, margin_right, y)
#             c.setStrokeColorRGB(0, 0, 0)  # Reset to black
#             y -= 15
    
#     # Add footer
#     c.setFont("Helvetica", 8)
#     c.drawString(margin_left, 30, f"Total Recommendations: {len(recommendations)}")
#     c.drawRightString(margin_right, 30, "End of Report")
    
#     c.save()
#     buffer.seek(0)
#     return buffer

def generate_pdf(recommendations):
    """Advanced PDF generation using ReportLab's Platypus framework"""
    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter, topMargin=1*inch, bottomMargin=1*inch)
    
    # Get styles
    styles = getSampleStyleSheet()
    title_style = styles['Title']
    heading_style = styles['Heading2']
    normal_style = styles['Normal']
    
    # Build story (content)
    story = []
    
    # Title
    title = Paragraph("Course Recommendations Report", title_style)
    story.append(title)
    story.append(Spacer(1, 0.3*inch))
    
    # Date
    from datetime import datetime
    date_text = f"Generated on: {datetime.now().strftime('%B %d, %Y at %I:%M %p')}"
    date_para = Paragraph(date_text, normal_style)
    story.append(date_para)
    story.append(Spacer(1, 0.2*inch))
    
    for i, course in enumerate(recommendations, 1):
        # Extract course information
        code = course.get('course_code', 'N/A')
        name = course.get('course_name', 'Untitled Course')
        desc = course.get('description', 'No description available.')
        outcomes = course.get('learning_outcomes', []) or course.get('outcomes', [])
        
        # Handle similarity
        similarity = course.get('adjusted_similarity', 0)
        try:
            sim_percentage = convert_similarity_to_percentage(similarity) if 'convert_similarity_to_percentage' in globals() else f"{float(similarity)*100:.1f}"
            sim_text = f"{sim_percentage}%"
        except:
            sim_text = "N/A"
        
        # Course header
        header_text = f"{i}. {code} - {name} (Match: {sim_text})"
        header_para = Paragraph(header_text, heading_style)
        story.append(header_para)
        story.append(Spacer(1, 0.1*inch))
        
        # Description
        desc_text = f"<b>Description:</b> {desc}"
        desc_para = Paragraph(desc_text, normal_style)
        story.append(desc_para)
        story.append(Spacer(1, 0.1*inch))
        
        # Learning Outcomes
        if outcomes and len(outcomes) > 0:
            outcomes_text = "<b>Learning Outcomes:</b><br/>"
            for outcome in outcomes:
                outcome_clean = str(outcome).strip()
                if outcome_clean:
                    outcomes_text += f"• {outcome_clean}<br/>"
            
            outcomes_para = Paragraph(outcomes_text, normal_style)
            story.append(outcomes_para)
        
        story.append(Spacer(1, 0.3*inch))
    
    # Build PDF
    doc.build(story)
    buffer.seek(0)
    return buffer

def load_embeddings_data(filepath):
    """Load embeddings JSON and normalize into a list of dicts."""
    try:
        with open(filepath, 'r', encoding='utf-8') as f:
            data = json.load(f)

        # 1) Extract the raw list
        if isinstance(data, dict) and "embeddings" in data:
            raw = data["embeddings"]
        elif isinstance(data, list):
            raw = data
        else:
            st.error(f"Unexpected format: top‐level is {type(data)}")
            return []

        # 2) Normalize into a list of dicts
        items = []
        for entry in raw:
            if isinstance(entry, dict):
                items.append(entry)
            else:
                # skip anything that isn't a dict
                st.warning(f"Skipping non‐dict entry of type {type(entry)}")
        
        return items

    except FileNotFoundError:
        st.error(f"Embeddings file not found: {filepath}")
        return []
    except json.JSONDecodeError as e:
        st.error(f"Error parsing embeddings file: {e}")
        return []
    except Exception as e:
        st.error(f"Error loading embeddings: {str(e)}")
        return []

courses = load_embeddings_data('all_course_embeddings.json')

# Additional helper function for debugging
def show_course_type_distribution(courses):
    """Show distribution of course types for debugging"""
    if not courses:
        return
    
    type_counts = {'technical': 0, 'business': 0, 'mixed': 0}
    for course in courses:
        course_type = course.get('course_type', 'unknown')
        if course_type in type_counts:
            type_counts[course_type] += 1
    
    st.write("**Course Type Distribution:**")
    for course_type, count in type_counts.items():
        st.write(f"- {course_type.title()}: {count}")


def next_step(): 
    st.session_state.current_step += 1

def prev_step(): 
    st.session_state.current_step -= 1

def generate_smart_question(requirements_data, conversation_history, question_count, 
                          current_tech_business_ratio=None):
    """
    Generate smart follow-up questions, including tech/business ratio clarification
    """
    try:
        if st.session_state.openai_api_key:
            client = OpenAI(api_key=st.session_state.openai_api_key)
            context = f"Initial Requirements: {st.session_state.initial_req}\n\n"
            
            if current_tech_business_ratio is not None:
                context += f"Current Tech/Business Ratio: {current_tech_business_ratio} "
                context += f"({'Technical focus' if current_tech_business_ratio > 0.6 else 'Business focus' if current_tech_business_ratio < 0.4 else 'Balanced'})\n\n"
            
            if conversation_history:
                context += "Previous conversation:\n"
                for msg in conversation_history[-4:]:
                    if msg['type'] == 'user':
                        context += f"User response: {msg['content']}\n"
            
            prompt = f"""You are a corporate training consultant. Based on the conversation, generate ONE specific follow-up question to understand training needs better.

Focus on areas like:
- Technical vs business skill balance
- Specific skill gaps
- Training objectives
- Timeline and constraints
- Success metrics
- Team composition

{context}

Generate only the question:"""

            response = client.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=[{"role": "user", "content": prompt}],
                max_tokens=100, temperature=0.7
            )
            return response.choices[0].message.content.strip()
        else:
            raise Exception("No API key")
    except:
        fallback_questions = [
            "What's the ideal balance between technical and business skills for your team?",
            "What specific technical or business challenges are you facing?",
            "How would you measure the success of this training program?",
            "What's your timeline for implementing this training?",
            "Which skills would have the most immediate impact on your team's performance?",
            "Are there any specific tools or technologies your team needs to learn?"
        ]
        return fallback_questions[min(question_count, len(fallback_questions)-1)]


# MongoDB connection
# client = MongoClient("mongodb+srv://mkunchal:gvL1yK8iNrIL3DX3@otpilearn.hjnl2ii.mongodb.net/")
# db = client.otpilearn
# collection = db.company

def extract_key_insights(text):
    insights = {'industry_focus': [], 'training_type': [], 'participant_count': '', 'timeline': '', 'delivery_method': [], 'objectives': []}
    text_lower = text.lower()
    
    for industry in ['technology', 'healthcare', 'finance', 'manufacturing', 'education', 'retail']:
        if industry in text_lower: insights['industry_focus'].append(industry)
    
    for t_type in ['leadership', 'technical', 'compliance', 'soft skills', 'project management']:
        if t_type.replace(' ', '') in text_lower.replace(' ', ''): insights['training_type'].append(t_type)
    
    numbers = re.findall(r'\d+', text)
    if numbers: insights['participant_count'] = numbers[0]
    
    return insights

def update_requirements_data(requirements_data, new_insights):
    for key, value in new_insights.items():
        if isinstance(value, list):
            if key not in requirements_data: requirements_data[key] = []
            requirements_data[key].extend(value)
            requirements_data[key] = list(set(requirements_data[key]))
        elif value and value.strip():
            requirements_data[key] = value
    return requirements_data
    

def get_text_embedding(text, client):
    """Generate embedding for user requirements text"""
    try:
        response = client.embeddings.create(
            model="text-embedding-3-small",
            input=text
        )
        return response.data[0].embedding
    except Exception as e:
        st.error(f"Error generating embedding: {str(e)}")
        return None
    
def classify_course_type(course_data):
    """
    Classify course as technical or business based on course content
    Returns: 'technical', 'business', or 'mixed'
    """
    # Keywords for technical courses
    technical_keywords = [
        'programming', 'algorithm', 'data structure', 'machine learning', 'ai', 'artificial intelligence',
        'database', 'software', 'computer', 'coding', 'python', 'java', 'javascript', 'web development',
        'cybersecurity', 'network', 'cloud', 'devops', 'api', 'backend', 'frontend', 'mobile app',
        'data science', 'analytics', 'statistics', 'mathematical', 'engineering', 'technical',
        'system design', 'architecture', 'framework', 'library', 'debugging', 'testing'
    ]
    
    # Keywords for business courses
    business_keywords = [
        'management', 'leadership', 'strategy', 'marketing', 'finance', 'accounting', 'business',
        'entrepreneurship', 'operations', 'project management', 'communication', 'negotiation',
        'sales', 'customer', 'organization', 'team building', 'hr', 'human resources',
        'economics', 'consulting', 'planning', 'decision making', 'process improvement',
        'change management', 'innovation', 'corporate', 'executive', 'administration'
    ]
    
    # Combine course content for analysis
    content = f"{course_data.get('course_name', '')} {course_data.get('course_description', '')} {course_data.get('embedding_text', '')}"
    content = content.lower()
    
    # Count keyword matches
    tech_count = sum(1 for keyword in technical_keywords if keyword in content)
    business_count = sum(1 for keyword in business_keywords if keyword in content)
    
    # Classification logic
    if tech_count > business_count * 1.5:
        return 'technical'
    elif business_count > tech_count * 1.5:
        return 'business'
    else:
        return 'mixed'
    
def calculate_similarity_with_requirements(user_embedding, course_embedding):
    """Calculate cosine similarity between user requirements and course"""
    try:
        if not user_embedding or not course_embedding:
            return 0.0
        
        # Convert to numpy arrays
        user_vec = np.array(user_embedding).reshape(1, -1)
        course_vec = np.array(course_embedding).reshape(1, -1)
        
        # Calculate cosine similarity
        similarity = cosine_similarity(user_vec, course_vec)[0][0]
        return float(similarity)
    except Exception as e:
        print(f"Error calculating similarity: {e}")
        return 0.0


def convert_similarity_to_percentage(similarity_score):
    """Convert similarity score to percentage for display"""
    # Raw cosine similarity is typically between -1 and 1
    # Convert to 0-100% scale
    if similarity_score < 0:
        return 0
    return int(similarity_score * 100)

def filter_by_tech_business_ratio(courses, tech_ratio):
    """
    Filter courses based on user's technical vs business preference
    tech_ratio: 0-100 (percentage for technical focus)
    """
    business_ratio = 100 - tech_ratio
    filtered_courses = []
    
    for course in courses:
        course_type = classify_course_type(course)
        
        # Scoring logic based on user preference
        if tech_ratio >= 70:  # Heavy technical focus
            if course_type == 'technical':
                course['type_match_score'] = 1.0
            elif course_type == 'mixed':
                course['type_match_score'] = 0.7
            else:  # business
                course['type_match_score'] = 0.3
        elif tech_ratio <= 30:  # Heavy business focus
            if course_type == 'business':
                course['type_match_score'] = 1.0
            elif course_type == 'mixed':
                course['type_match_score'] = 0.7
            else:  # technical
                course['type_match_score'] = 0.3
        else:  # Balanced approach (30-70%)
            if course_type == 'mixed':
                course['type_match_score'] = 1.0
            else:
                course['type_match_score'] = 0.8
        
        # Adjust final similarity score based on type match
        course['adjusted_similarity'] = course['similarity_score'] * course['type_match_score']
        course['course_type'] = course_type
        
        filtered_courses.append(course)
    
    return filtered_courses

def generate_course_recommendations(user_requirements_text):
    """
    Generate course recommendations based on user requirements using cosine similarity,
    then dedupe modules so each course_code appears only once (highest similarity).
    """
    try:
        # 1) Load embeddings
        embeddings_data = load_embeddings_data('all_course_embeddings.json')
        if not embeddings_data:
            st.error("No course embeddings data found!")
            return []

        # 2) Ensure API key & get client
        if not st.session_state.openai_api_key:
            st.error("OpenAI API key not found!")
            return []
        client = OpenAI(api_key=st.session_state.openai_api_key)

        # 3) Embed user requirements
        user_embedding = get_text_embedding(user_requirements_text, client)
        if not user_embedding:
            st.error("Failed to generate embedding for user requirements!")
            return []

        # 4) Compute module-level similarities
        module_sims = []
        for course_data in embeddings_data:
            vector = course_data.get('vector')
            if not vector:
                continue
            sim = calculate_similarity_with_requirements(user_embedding, vector)
            if sim >= 0.2:  # your threshold (e.g. 0.4 for 40%)
                module_sims.append({
                    'course_code':      course_data.get('course_code', 'N/A'),
                    'course_name':      course_data.get('course_name', 'N/A'),
                    'description':      course_data.get('course_description', ''),
                    'similarity_score': sim,
                    'best_week':        course_data.get('module_display_name', 'overview'),
                    'module_title':     ', '.join(course_data.get('topics', [])),
                    'learning_outcomes': course_data.get('learning_outcomes', []),
                    'raw_data':         course_data
                })

        # 5) Dedupe by course_code—keep only the module with max similarity
        best_per_course = {}
        for info in module_sims:
            code = info['course_code']
            if code not in best_per_course or info['similarity_score'] > best_per_course[code]['similarity_score']:
                best_per_course[code] = info
        course_similarities = list(best_per_course.values())

        # 6) Filter by tech/business preference
        tech_ratio    = st.session_state.form_data.get('tech_ratio', 50)
        filtered      = filter_by_tech_business_ratio(course_similarities, tech_ratio)

        # 7) Sort & rank
        filtered.sort(key=lambda x: x['adjusted_similarity'], reverse=True)
        top_recs = filtered[:10]
        for idx, course in enumerate(top_recs, start=1):
            course['rank'] = idx

        # 8) Debug stats
        st.write("📊 **Recommendation Stats:**")
        # st.write(f"- Total modules in DB: {len(embeddings_data)}")
        # st.write(f"- Modules above threshold: {len(module_sims)}")
        st.write(f"- Unique courses after dedupe: {len(course_similarities)}")
        st.write(f"- Final recommendations: {len(top_recs)}")
        st.write(f"- Tech/Business ratio: {tech_ratio}%/{100-tech_ratio}%")
        if top_recs:
            top = top_recs[0]
            pct = convert_similarity_to_percentage(top['adjusted_similarity'])
            st.write(f"- Top match: {top['course_name']} ({pct}%)")

        return top_recs

    except Exception as e:
        st.error(f"Error generating recommendations: {e}")
        return []

    

# Header
st.markdown("""
<div class="header-section">
    <h1 class="header-title">Stevens OptiLearn</h1>
    <p class="header-subtitle">Intelligent Corporate Training Solutions powered by Stevens Institute of Technology</p>
</div>
""", unsafe_allow_html=True)

st.markdown(
    """
    <style>
    /* target all primary buttons – if this is your only primary button it will apply just to it */
    div.stButton > button[kind="primary"] {
        color: white !important;
    }
    </style>
    """,
    unsafe_allow_html=True,
)

# Login System
if not st.session_state.logged_in:
    st.markdown('<div class="form-wrapper"><div class="section-title" style="text-align: center;">Corporate Access Portal</div></div>', unsafe_allow_html=True)
    
    col1, col2, col3 = st.columns([1, 2, 1])
    with col2:
        user_id = st.text_input("User ID", placeholder="Enter your Stevens corporate ID")
        password = st.text_input("Password", type="password", placeholder="Enter your password")
        api_key = "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"
        
        if st.button("Login"):
            if user_id in USER_CREDENTIALS and USER_CREDENTIALS[user_id] == password:
                st.session_state.logged_in = True
                st.session_state.current_step = 2
                if api_key: st.session_state.openai_api_key = api_key
                st.success("Login successful!")
                st.rerun()
            else:
                st.error("Invalid credentials. Please try again.")
    st.stop()

# Progress Indicator
if st.session_state.logged_in:
    steps = ["Company Info", "Training Focus", "Requirements", "AI Consultation", "Recommendations"]
    progress_html = '<div class="progress-indicator">'
    for i, step in enumerate(steps, 2):
        status = "completed" if st.session_state.current_step > i else "active" if st.session_state.current_step == i else ""
        progress_html += f'<div class="progress-step {status}">{step}</div>'
    progress_html += '</div>'
    st.markdown(progress_html, unsafe_allow_html=True)

# Step 2: Company Information
if st.session_state.logged_in and st.session_state.current_step == 2:
    st.markdown('<div class="form-wrapper"><div class="section-title">Company Information</div>', unsafe_allow_html=True)
    
    col1, col2 = st.columns(2)
    with col1:
        company_name = st.text_input("Company Name *", value=st.session_state.form_data.get("company_name", ""))
        industry = st.selectbox("Industry *", ["", "Technology", "Finance", "Healthcare", "Manufacturing", "Retail", "Consulting", "Education", "Government", "Other"], 
                              index=["", "Technology", "Finance", "Healthcare", "Manufacturing", "Retail", "Consulting", "Education", "Government", "Other"].index(st.session_state.form_data.get("industry", "")))
    
    with col2:
        company_size = st.selectbox("Company Size *", ["", "1-50 employees", "51-200 employees", "201-1000 employees", "1001-5000 employees", "5000+ employees"],
                                  index=["", "1-50 employees", "51-200 employees", "201-1000 employees", "1001-5000 employees", "5000+ employees"].index(st.session_state.form_data.get("company_size", "")))
        contact_email = st.text_input("Contact Email *", value=st.session_state.form_data.get("contact_email", ""))
        
        # Email validation
        if contact_email and not is_valid_email(contact_email):
            st.markdown('<div class="email-error">⚠️ Please enter a valid email address</div>', unsafe_allow_html=True)
    
    if st.button("Continue →"):
        if all([company_name, industry, company_size, contact_email]):
            if is_valid_email(contact_email):
                st.session_state.form_data.update({"company_name": company_name, "industry": industry, "company_size": company_size, "contact_email": contact_email})
                company={"company_name": company_name, "industry": industry, "company_size": company_size, "contact_email": contact_email}
                #collection.insert_one(company)
                st.success("Company information saved successfully!")
                next_step()
                st.rerun()
            else:
                st.error("Please enter a valid email address.")
        else:
            st.error("Please fill in all required fields.")
    st.markdown('</div>', unsafe_allow_html=True)

# Step 3: Training Focus (Tech vs Business Ratio)
elif st.session_state.logged_in and st.session_state.current_step == 3:
    st.markdown('<div class="form-wrapper"><div class="section-title">Training Focus Areas</div>', unsafe_allow_html=True)
    
    st.markdown('<div style="text-align: center; margin: 2rem 0;"><h3 style="font-family: \'Roboto Slab\', serif; color: #2c3e50;">Technical vs. Business Skills Distribution</h3></div>', unsafe_allow_html=True)
    
    # Fixed slider with proper label
    tech_ratio = st.slider(
        "Technical Skills Percentage", 
        min_value=0, 
        max_value=100, 
        value=st.session_state.form_data.get("tech_ratio", 50),
        help="Slide to adjust the balance between technical and business skills",
        label_visibility="hidden"
    )
    business_ratio = 100 - tech_ratio
    
    st.markdown(f"""
    <div class="ratio-display">
        <div class="ratio-item">
            <div class="ratio-value">{business_ratio}%</div>
            <div class="ratio-label">Business Skills</div>
        </div>
        <div class="ratio-item">
            <div class="ratio-value">{tech_ratio}%</div>
            <div class="ratio-label">Technical Skills</div>
        </div>
    </div>
    """, unsafe_allow_html=True)
   
    
    # Update session state immediately when selection changes
    st.session_state.form_data.update({
        "tech_ratio": tech_ratio, 
        "business_ratio": business_ratio, 
    })
    

    
    # Navigation buttons
    col1, col2 = st.columns(2)
    with col1:
        if st.button("← Back", key="step3_back"): 
            prev_step()
            st.rerun()
    with col2:
        if st.button("Continue →", key="step3_continue"): 
            # Ensure data is saved before moving forward
            st.session_state.form_data.update({
                "tech_ratio": tech_ratio, 
                "business_ratio": business_ratio, 
            })
            next_step()
            st.rerun()
    
    st.markdown('</div>', unsafe_allow_html=True)
    
# Step 4: Requirements Gathering
elif st.session_state.logged_in and st.session_state.current_step == 4:
    st.markdown('<div class="form-wrapper"><div class="section-title">Training Requirements</div>', unsafe_allow_html=True)
    
    # Add custom CSS for the burgundy button styling
    st.markdown("""
    <style>
    .upload-option {
        background: linear-gradient(135deg, #8B1538 0%, #A91B47 100%);
        color: white;
        padding: 1rem;
        border-radius: 10px;
        margin: 0.5rem 0;
        text-align: center;
        cursor: pointer;
        transition: all 0.3s ease;
        border: 2px solid transparent;
    }
    .upload-option:hover {
        background: linear-gradient(135deg, #A91B47 0%, #8B1538 100%);
        border-color: #D4AF37;
        transform: translateY(-2px);
        box-shadow: 0 4px 15px rgba(139, 21, 56, 0.3);
    }
    .upload-option.active {
        background: linear-gradient(135deg, #2C2C2C 0%, #404040 100%);
        border-color: #D4AF37;
    }
    .option-title {
        font-weight: bold;
        font-size: 1.1rem;
        margin-bottom: 0.5rem;
    }
    .option-description {
        font-size: 0.9rem;
        opacity: 0.9;
    }
    .divider-text {
        text-align: center;
        margin: 1.5rem 0;
        position: relative;
        color: #666;
    }
    .divider-text::before {
        content: '';
        position: absolute;
        top: 50%;
        left: 0;
        right: 0;
        height: 1px;
        background: linear-gradient(90deg, transparent, #666, transparent);
    }
    .divider-text span {
        background: white;
        padding: 0 1rem;
        font-weight: 500;
    }
    
    /* Custom burgundy button styling */
    .stButton > button {
        background: linear-gradient(135deg, #8B1538 0%, #A91B47 100%) !important;
        color: white !important;
        border: none !important;
        border-radius: 8px !important;
        padding: 0.75rem 1.5rem !important;
        font-weight: 600 !important;
        transition: all 0.3s ease !important;
        box-shadow: 0 2px 8px rgba(139, 21, 56, 0.2) !important;
    }
    
    .stButton > button:hover {
        background: linear-gradient(135deg, #A91B47 0%, #8B1538 100%) !important;
        transform: translateY(-2px) !important;
        box-shadow: 0 4px 15px rgba(139, 21, 56, 0.4) !important;
    }
    
    .stButton > button:active {
        transform: translateY(0) !important;
    }
    
    /* Selected button styling */
    .selected-option {
        background: linear-gradient(135deg, #6B1329 0%, #8B1538 100%) !important;
        border: 2px solid #D4AF37 !important;
    }
    </style>
    """, unsafe_allow_html=True)
    
    # Initialize input method selection
    if 'input_method' not in st.session_state:
        st.session_state.input_method = None
    
    # Input method selection with burgundy buttons
    st.markdown("**Choose your preferred input method:**")
    st.markdown("<br>", unsafe_allow_html=True)
    
    col1, col2 = st.columns(2)
    
    with col1:
        if st.button("📝 Write Text Description", 
                    key="text_option", 
                    help="Describe your requirements in text",
                    use_container_width=True):
            st.session_state.input_method = 'text'
            st.rerun()
    
    with col2:
        if st.button("📄 Upload Document", 
                    key="file_option", 
                    help="Upload PDF, DOC, or PPT file",
                    use_container_width=True):
            st.session_state.input_method = 'file'
            st.rerun()

    st.markdown("<br>", unsafe_allow_html=True)

    # Show content based on selected method
    requirements_text = ""
    
    # Text input method
    if st.session_state.input_method == 'text':
        st.markdown("### ✍️ Text Input")
        initial_req = st.text_area(
            "Describe your training requirements:", 
            height=200, 
            placeholder="Tell us about your training goals, target audience, specific skills needed, preferred delivery format, timeline, and desired outcomes...",
            value=st.session_state.get('initial_req', '') or ""
        )
        requirements_text = initial_req
    
    # File upload method
    elif st.session_state.input_method == 'file':
        st.markdown("### 📄 Document Upload")
        uploaded_file = st.file_uploader(
            "Choose a file", 
            type=['pdf', 'doc', 'docx', 'ppt', 'pptx', 'txt'],
            help="Supported formats: PDF, DOC, DOCX, PPT, PPTX, TXT"
        )
        
        if uploaded_file is not None:
            try:
                with st.spinner("📖 Extracting text from your document..."):
                    # Extract text based on file type
                    if uploaded_file.type == "application/pdf":
                        try:
                            # Try multiple PDF libraries for better extraction
                            requirements_text = ""
                            
                            # Method 1: Try PyPDF2 first
                            try:
                                import PyPDF2
                                pdf_reader = PyPDF2.PdfReader(uploaded_file)
                                text_content = ""
                                for page in pdf_reader.pages:
                                    page_text = page.extract_text()
                                    if page_text:
                                        text_content += page_text + "\n"
                                requirements_text = text_content.strip()
                                
                                if not requirements_text:
                                    raise Exception("No text extracted with PyPDF2")
                                    
                            except (ImportError, Exception) as e:
                                # Method 2: Try pdfplumber as fallback
                                try:
                                    import pdfplumber
                                    uploaded_file.seek(0)  # Reset file pointer
                                    with pdfplumber.open(uploaded_file) as pdf:
                                        text_content = ""
                                        for page in pdf.pages:
                                            page_text = page.extract_text()
                                            if page_text:
                                                text_content += page_text + "\n"
                                        requirements_text = text_content.strip()
                                        
                                    if not requirements_text:
                                        raise Exception("No text extracted with pdfplumber")
                                        
                                except (ImportError, Exception):
                                    # Method 3: Try pymupdf (fitz) as final fallback
                                    try:
                                        import fitz  # PyMuPDF
                                        uploaded_file.seek(0)  # Reset file pointer
                                        pdf_document = fitz.open(stream=uploaded_file.read(), filetype="pdf")
                                        text_content = ""
                                        for page_num in range(pdf_document.page_count):
                                            page = pdf_document[page_num]
                                            page_text = page.get_text()
                                            if page_text:
                                                text_content += page_text + "\n"
                                        pdf_document.close()
                                        requirements_text = text_content.strip()
                                        
                                        if not requirements_text:
                                            raise Exception("No text extracted with PyMuPDF")
                                            
                                    except (ImportError, Exception):
                                        st.error("❌ Could not extract text from PDF. This might be a scanned document or have complex formatting.")
                                        st.info("💡 Try one of these alternatives:")
                                        st.info("• Copy and paste the text manually using the 'Write Text Description' option")
                                        st.info("• Convert the PDF to a Word document first")
                                        st.info("• Use a different PDF file")
                                        requirements_text = ""
                                        
                        except Exception as e:
                            st.error(f"❌ Error processing PDF: {str(e)}")
                            requirements_text = ""
                    
                    elif uploaded_file.type in ["application/vnd.openxmlformats-officedocument.wordprocessingml.document", 
                                              "application/msword"]:
                        try:
                            import docx
                            doc = docx.Document(uploaded_file)
                            text_content = ""
                            for paragraph in doc.paragraphs:
                                text_content += paragraph.text + "\n"
                            requirements_text = text_content.strip()
                        except ImportError:
                            st.error("python-docx library not available. Please install it or use text input.")
                            requirements_text = ""
                    
                    elif uploaded_file.type in ["application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                              "application/vnd.ms-powerpoint"]:
                        try:
                            from pptx import Presentation
                            prs = Presentation(uploaded_file)
                            text_content = ""
                            for slide in prs.slides:
                                for shape in slide.shapes:
                                    if hasattr(shape, "text"):
                                        text_content += shape.text + "\n"
                            requirements_text = text_content.strip()
                        except ImportError:
                            st.error("python-pptx library not available. Please install it or use text input.")
                            requirements_text = ""
                    
                    elif uploaded_file.type == "text/plain":
                        # Handle text files
                        requirements_text = str(uploaded_file.read(), "utf-8")
                    
                    else:
                        st.error("Unsupported file type. Please upload PDF, DOC, DOCX, PPT, PPTX, or TXT files.")
                        requirements_text = ""
                
                if requirements_text:
                    st.success(f"✅ Successfully extracted text from {uploaded_file.name}")
                    
                    # Show preview of extracted text
                    with st.expander("📋 Preview Extracted Content", expanded=False):
                        preview_text = requirements_text[:1000] + ("..." if len(requirements_text) > 1000 else "")
                        st.text_area("Extracted Text:", value=preview_text, height=150, disabled=True)
                        if len(requirements_text) > 1000:
                            st.info(f"Showing first 1000 characters of {len(requirements_text)} total characters")
                    
                    # Allow editing of extracted text
                    st.markdown("**Review and edit the extracted text if needed:**")
                    requirements_text = st.text_area("Edit requirements:", value=requirements_text, height=150)
                else:
                    st.error("❌ Could not extract text from the uploaded file. Please try a different file or use text input.")
                    
            except Exception as e:
                st.error(f"❌ Error processing file: {str(e)}")
                st.info("💡 Please try uploading a different file or use the text input option.")
        
        else:
            st.info("👆 Please upload a document containing your training requirements")
    
    else:
        st.info("👆 Please select an input method above to get started")
    
    st.markdown("<br><br>", unsafe_allow_html=True)
    
    # Navigation buttons
    col1, col2, col3 = st.columns([1, 2, 1])
    
    with col1:
        if st.button("← Back", use_container_width=True): 
            prev_step()
            st.rerun()
    
    with col3:
        if st.button("Start AI Consultation", use_container_width=True):
            if requirements_text and requirements_text.strip():
                st.session_state.initial_req = requirements_text
                
                # Add to conversation history if not already there
                if 'conversation_history' not in st.session_state:
                    st.session_state.conversation_history = []
                
                st.session_state.conversation_history.append({
                    "type": "user", 
                    "content": requirements_text, 
                    "timestamp": datetime.now().isoformat()
                })
                
                # Extract insights and generate question
                try:
                    initial_insights = extract_key_insights(requirements_text)
                    if 'requirements_data' not in st.session_state:
                        st.session_state.requirements_data = {}
                    st.session_state.requirements_data = update_requirements_data(
                        st.session_state.requirements_data, 
                        initial_insights
                    )
                    
                    # Generate first AI question
                    if 'question_count' not in st.session_state:
                        st.session_state.question_count = 0
                        
                    st.session_state.current_ai_question = generate_smart_question(
                        st.session_state.requirements_data, 
                        st.session_state.conversation_history, 
                        st.session_state.question_count
                    )
                    
                    next_step()
                    st.rerun()
                    
                except Exception as e:
                    st.error(f"Error processing requirements: {str(e)}")
                    st.info("Please try again or contact support if the issue persists.")
            else:
                if st.session_state.input_method == 'text':
                    st.error("Please provide your training requirements in the text area above.")
                elif st.session_state.input_method == 'file':
                    st.error("Please upload a document or switch to text input to provide your requirements.")
                else:
                    st.error("Please select an input method and provide your requirements.")
    
    st.markdown('</div>', unsafe_allow_html=True)

# Step 5: AI Consultation
elif st.session_state.logged_in and st.session_state.current_step == 5:
    st.markdown('<div class="ai-chat-container"><div class="section-title" style="margin: 0 0 2rem 0; text-align: center;">AI Training Consultant</div>', unsafe_allow_html=True)
    
    # Show current AI question if user hasn't reached question limit
    if st.session_state.question_count < 7:
        if st.session_state.current_ai_question:
            st.markdown(f'<div class="ai-message"><div class="ai-label">Stevens AI Consultant</div>{st.session_state.current_ai_question}</div>', unsafe_allow_html=True)
        
        # Use unique key for text area to force refresh
        user_response = st.text_area("Your response:", height=80, placeholder="Please share your thoughts...", key=f"response_{st.session_state.question_count}")
        
        col1, col2, col3 = st.columns(3)
        with col1:
            if st.button("Submit Response"):
                if user_response.strip():
                    # Add user response to history
                    st.session_state.conversation_history.append({"type": "user", "content": user_response, "timestamp": datetime.now().isoformat()})
                    # Extract insights
                    insights = extract_key_insights(user_response)
                    st.session_state.requirements_data = update_requirements_data(st.session_state.requirements_data, insights)
                    # Increment question count
                    st.session_state.question_count += 1
                    # Generate next question if not at limit
                    if st.session_state.question_count < 7:
                        st.session_state.current_ai_question = generate_smart_question(st.session_state.requirements_data, st.session_state.conversation_history, st.session_state.question_count)
                    else:
                        st.session_state.current_ai_question = None
                    st.rerun()
        with col2:
            if st.button("Skip Question"): 
                st.session_state.conversation_history.append({"type": "user", "content": "[Question Skipped]", "timestamp": datetime.now().isoformat()})
                st.session_state.question_count += 1
                if st.session_state.question_count < 7:
                    st.session_state.current_ai_question = generate_smart_question(st.session_state.requirements_data, st.session_state.conversation_history, st.session_state.question_count)
                else:
                    st.session_state.current_ai_question = None
                st.rerun()
        with col3:
            if st.button("Finish Consultation"): 
                # Generate course recommendations before moving to next step
                full_requirements = st.session_state.initial_req
                for msg in st.session_state.conversation_history:
                    if msg["type"] == "user" and msg["content"] != "[Question Skipped]":
                        full_requirements += " " + msg["content"]
                
                with st.spinner("🔍 Analyzing requirements and generating personalized recommendations..."):
                    st.session_state.course_recommendations = generate_course_recommendations(full_requirements)
                
                next_step()
                st.rerun()
    else:
        st.markdown('<div class="ai-message"><div class="ai-label">Stevens AI Consultant</div>Thank you for providing detailed information about your training needs. Let\'s generate your personalized course recommendations.</div>', unsafe_allow_html=True)
        if st.button("Generate Recommendations"): 
            # Generate course recommendations
            full_requirements = st.session_state.initial_req
            for msg in st.session_state.conversation_history:
                if msg["type"] == "user" and msg["content"] != "[Question Skipped]":
                    full_requirements += " " + msg["content"]
            
            with st.spinner("🔍 Analyzing requirements and generating personalized recommendations..."):
                st.session_state.course_recommendations = generate_course_recommendations(full_requirements)
            
            next_step()
            st.rerun()
    
    st.markdown('</div>', unsafe_allow_html=True)


# Step 6: Recommendations UI
elif st.session_state.logged_in and st.session_state.current_step == 6:
    # Retrieve recommendations and ratios
    recommendations = st.session_state.get('course_recommendations', [])
    tech_ratio      = st.session_state['form_data'].get('tech_ratio', 50)
    business_ratio  = 100 - tech_ratio

 

    if recommendations:
        # Cards Wrapper
        st.markdown('<div class="recommendations-list">', unsafe_allow_html=True)

        for course in recommendations:
            rank = course.get('rank', '?')
            code = course.get('course_code', 'N/A')
            name = course.get('course_name', 'Untitled Course')
            desc = course.get('description', 'No description available.')
            sim  = convert_similarity_to_percentage(course.get('adjusted_similarity', 0))
            orig = convert_similarity_to_percentage(course.get('similarity_score', 0))
            ctype = course.get('course_type', 'mixed')
            type_labels = {'technical': 'Technical', 'business': 'Business', 'mixed': 'Hybrid'}
            type_icons  = {'technical': '💻', 'business': '📊', 'mixed': '⚖️'}
            label = type_labels.get(ctype, 'Hybrid')
            icon  = type_icons.get(ctype, '⚖️')

            st.markdown(f"""
            <article class="professional-course-card">
              <div class="course-header-pro">
                <div class="course-rank-pro">#{rank}</div>
                <div class="course-code-pro">{code}</div>
                <h3 class="course-name-pro">{name}</h3>
                <div class="course-type-badge">{icon} {label} Track</div>
              </div>
              <div class="course-body-pro">
                <div class="similarity-section-pro">
                  <div class="similarity-metric">
                    <div class="similarity-value-pro">{sim}%</div>
                    <div class="similarity-label-pro">Overall Match</div>
                  </div>
                  <div class="similarity-metric">
                    <div class="similarity-value-pro">{orig}%</div>
                    <div class="similarity-label-pro">Content Match</div>
                  </div>
                </div>
                <div class="course-description-pro">
                  <strong>Course Overview:</strong><br>{desc}
                </div>
              </div>
            </article>
            """, unsafe_allow_html=True)

        # close cards wrapper
        st.markdown('</div>', unsafe_allow_html=True)

        # Executive Summary
        # compute stats
        scores       = [convert_similarity_to_percentage(c.get('adjusted_similarity',0)) for c in recommendations]
        avg_score    = int(sum(scores)/len(scores)) if scores else 0
        tech_count   = sum(1 for c in recommendations if c.get('course_type')=='technical')
        biz_count    = sum(1 for c in recommendations if c.get('course_type')=='business')
        hybrid_count = sum(1 for c in recommendations if c.get('course_type')=='mixed')

        st.markdown(f"""
        <div class="recommendations-summary">
          <h2 style="text-align:center; font-family:'Roboto Slab', serif; color: white;">Executive Summary</h2>
          <div class="summary-grid">
            <div class="summary-metric">
              <div class="summary-metric-value">{len(recommendations)}</div>
              <div class="summary-metric-label">Total Courses</div>
            </div>
            <div class="summary-metric">
              <div class="summary-metric-value">{avg_score}%</div>
              <div class="summary-metric-label">Avg. Match Score</div>
            </div>
            <div class="summary-metric">
              <div class="summary-metric-value">{tech_count}</div>
              <div class="summary-metric-label">Technical Courses</div>
            </div>
            <div class="summary-metric">
              <div class="summary-metric-value">{biz_count}</div>
              <div class="summary-metric-label">Business Courses</div>
            </div>
            <div class="summary-metric">
              <div class="summary-metric-value">{hybrid_count}</div>
              <div class="summary-metric-label">Hybrid Courses</div>
            </div>
            <div class="summary-metric">
              <div class="summary-metric-value">{tech_ratio}%</div>
              <div class="summary-metric-label">Tech Focus</div>
            </div>
          </div>
        </div>
        """, unsafe_allow_html=True)


        # Next Steps
        #st.markdown("### Next Steps")
        #st.markdown("<p style='margin-bottom: 2rem; color: #6c757d;'>Choose how you'd like to proceed with your personalized training plan:</p>", unsafe_allow_html=True)
  
        col1, col2, col3 = st.columns(3)
    
        with col2:
            if st.button("Download PDF", key="Download PDF"):
                # reset to step 2
                pdf_buffer = generate_pdf(recommendations)
                st.download_button(
                label="📥 Download Course Recommendations as PDF",
                data=pdf_buffer,
                file_name="course_recommendations.pdf",
                mime="application/pdf"
)
        with col3:
            if st.button("🔄 Start New Search", key="restart"):
                # reset to step 2
                st.session_state.current_step = 2
                st.session_state.form_data = {}
                st.session_state.initial_req = None
                st.session_state.conversation_history = []
                st.session_state.question_count = 0
                st.session_state.requirements_data = {}
                st.session_state.course_recommendations = []
                st.session_state.current_ai_question = None
                st.rerun()

        st.markdown('</div>', unsafe_allow_html=True)

    else:
        # No Results
        st.markdown("""
        <div class="no-results">
          <h3>🔍 No Matching Courses Found</h3>
          <p>We couldn't find any courses that match your specific requirements.<br>
             This might be because your criteria are very specialized.</p>
          <p style="margin-top: 1.5rem;"><strong>What you can do:</strong></p>
          <ul style="text-align: left; display: inline-block; margin-top: 1rem;">
            <li>Adjust your technical/business ratio</li>
            <li>Broaden your requirements</li>
            <li>Contact us for custom course development</li>
          </ul>
          <div style="margin-top: 2rem;">
            <button onclick="window.location.reload();" style="background: #C8102E; color: white; padding: 0.75rem 2rem; border: none; border-radius: 25px; cursor: pointer; font-weight: 600;">
              Try Again
            </button>
          </div>
        </div>
        """, unsafe_allow_html=True)


# Footer
if st.session_state.get('logged_in', False):
    current_step = st.session_state.get('current_step', 1)
    st.markdown(f"""
    <div style="text-align: center; padding: 3rem 1rem 2rem; color: #6c757d; font-size: 0.95rem; border-top: 2px solid #e9ecef; margin-top: 4rem; background: #f8f9fa;">
        <p style="margin-bottom: 1rem;"><strong style="color: #800020; font-size: 1.1rem;">Stevens Institute of Technology</strong></p>
        <p style="margin-bottom: 0.5rem;">College of Professional Education</p>
        <p style="margin-bottom: 1.5rem;">
            📧 cpe@stevens.edu &nbsp;|&nbsp; 📞 (201) 216-5000 &nbsp;|&nbsp; 🌐 stevens.edu/
        </p>
        <p style="font-size: 0.85rem; color: #adb5bd;">
            Session ID: SIT_{datetime.now().strftime('%Y%m%d_%H%M%S')} &nbsp;|&nbsp; Progress: Step {current_step-1} of 5
        </p>
    </div>
    """, unsafe_allow_html=True)